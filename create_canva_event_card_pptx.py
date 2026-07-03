from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "canva-event-card-layers.pptx"


def add_text(slide, text, left, top, width, height, font_size, bold=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    return box


prs = Presentation()
prs.slide_width = Inches(8.2)
prs.slide_height = Inches(1.0)

slide = prs.slides.add_slide(prs.slide_layouts[6])

scale = prs.slide_width / 492


def sx(px):
    return int(px * scale)


def sy(px):
    return int(px * scale)


card = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    sx(0),
    sy(0),
    sx(492),
    sy(59),
)
card.name = "卡片背景"
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(228, 154, 130)
card.line.fill.background()

badge = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    sx(14),
    sy(5),
    sx(52),
    sy(49),
)
badge.name = "服務勳章外框"
badge.fill.background()
badge.line.color.rgb = RGBColor(255, 255, 255)
badge.line.width = Pt(3)

add_text(slide, "剪", sx(26), sy(13), sx(40), sy(34), 30).name = "服務文字-剪"
add_text(slide, "林*2", sx(95), sy(19), sx(80), sy(28), 23).name = "顧客姓名"
add_text(slide, "0920-381348", sx(311), sy(17), sx(135), sy(24), 20).name = "電話"
add_text(slide, "婷", sx(470), sy(21), sx(20), sy(18), 14).name = "設計師"

prs.save(OUT)
print(OUT)
